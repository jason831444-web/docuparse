from pathlib import Path
from typing import Any


class OfficeExtractionService:
    def extract(self, path: Path, extension: str) -> tuple[str, list[dict[str, Any]], list[str], dict[str, Any]]:
        if extension == "docx":
            return self._docx(path)
        if extension == "xlsx":
            return self._xlsx(path)
        if extension == "pptx":
            return self._pptx(path)
        return "", [], [f"{extension.upper()} is not directly supported by the Office extractor."], {}

    def _docx(self, path: Path) -> tuple[str, list[dict[str, Any]], list[str], dict[str, Any]]:
        warnings: list[str] = []
        try:
            from docx import Document as DocxDocument
        except Exception as exc:
            return "", [], [f"DOCX extraction dependency is unavailable: {exc}."], {}

        document = DocxDocument(path)
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        table_lines = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_lines.append(" | ".join(cells))
        lines = paragraphs + table_lines
        metadata = {"paragraph_count": len(paragraphs), "table_count": len(document.tables)}
        blocks = [{"type": "paragraph", "content": value} for value in paragraphs[:200]]
        blocks.extend({"type": "table_row", "content": value} for value in table_lines[:200])
        return "\n".join(lines), blocks, warnings, metadata

    def _xlsx(self, path: Path) -> tuple[str, list[dict[str, Any]], list[str], dict[str, Any]]:
        warnings: list[str] = []
        try:
            from openpyxl import load_workbook
        except Exception as exc:
            return "", [], [f"XLSX extraction dependency is unavailable: {exc}."], {}

        workbook = load_workbook(path, read_only=True, data_only=True)
        lines = []
        blocks = []
        for sheet in workbook.worksheets:
            lines.append(f"Sheet: {sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                cells = [self._cell_text(cell) for cell in row]
                cells = [cell for cell in cells if cell]
                if not cells:
                    continue
                row_count += 1
                if row_count > 250:
                    warnings.append(f"Sheet '{sheet.title}' extraction limited to 250 non-empty rows.")
                    break
                line = " | ".join(cells)
                lines.append(line)
                if len(blocks) < 300:
                    blocks.append({"type": "sheet_row", "sheet": sheet.title, "content": line})
        metadata = {"sheet_count": len(workbook.worksheets), "sheet_names": [sheet.title for sheet in workbook.worksheets]}
        workbook.close()
        return "\n".join(lines), blocks, warnings, metadata

    def _pptx(self, path: Path) -> tuple[str, list[dict[str, Any]], list[str], dict[str, Any]]:
        warnings: list[str] = []
        try:
            from pptx import Presentation
        except Exception as exc:
            return "", [], [f"PPTX extraction dependency is unavailable: {exc}."], {}

        presentation = Presentation(path)
        lines = []
        blocks = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    for line in shape.text.splitlines():
                        cleaned = line.strip()
                        if cleaned:
                            slide_lines.append(cleaned)
            if slide_lines:
                lines.append(f"Slide {index}")
                lines.extend(slide_lines)
                blocks.append({"type": "slide", "slide": index, "content": "\n".join(slide_lines)})
        metadata = {"slide_count": len(presentation.slides)}
        return "\n".join(lines), blocks, warnings, metadata

    def _cell_text(self, value: Any) -> str:
        if value is None:
            return ""
        return str(value).strip()
